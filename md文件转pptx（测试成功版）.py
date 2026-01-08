#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
大纲转PPT工具 v1.1 (修复版)
修复段落格式设置问题
"""

import sys
import re
import os
import subprocess

from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QTextEdit, QLabel, QLineEdit, QComboBox, QSpinBox, QPushButton,
    QFileDialog, QMessageBox, QGroupBox, QFormLayout, QCheckBox,
    QStatusBar, QToolBar, QFrame, QDoubleSpinBox
)
from PyQt6.QtCore import Qt, QSettings
from PyQt6.QtGui import QFont, QAction, QKeySequence, QDragEnterEvent, QDropEvent

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsmap
from pptx.enum.text import PP_ALIGN
from lxml import etree


# ==================== 配色主题 ====================
THEMES = {
    "经典蓝": {"title_color": (0, 51, 102), "body_color": (51, 51, 51)},
    "商务灰": {"title_color": (64, 64, 64), "body_color": (89, 89, 89)},
    "活力橙": {"title_color": (204, 85, 0), "body_color": (51, 51, 51)},
    "清新绿": {"title_color": (0, 102, 51), "body_color": (51, 51, 51)},
    "优雅紫": {"title_color": (75, 0, 130), "body_color": (51, 51, 51)},
    "纯黑白": {"title_color": (0, 0, 0), "body_color": (33, 33, 33)},
}


def get_rgb_color(color_tuple):
    return RGBColor(color_tuple[0], color_tuple[1], color_tuple[2])


class DragDropTextEdit(QTextEdit):
    """支持拖拽的文本框"""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setAcceptDrops(True)

    def dragEnterEvent(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
        else:
            super().dragEnterEvent(event)

    def dropEvent(self, event: QDropEvent):
        if event.mimeData().hasUrls():
            for url in event.mimeData().urls():
                file_path = url.toLocalFile()
                if file_path.lower().endswith(('.txt', '.md', '.markdown')):
                    try:
                        content = None
                        for enc in ['utf-8', 'gbk', 'gb2312', 'utf-16']:
                            try:
                                with open(file_path, 'r', encoding=enc) as f:
                                    content = f.read()
                                break
                            except UnicodeDecodeError:
                                continue
                        if content:
                            self.setPlainText(content)
                            if hasattr(self.window(), 'status_bar'):
                                self.window().status_bar.showMessage(f"已导入: {os.path.basename(file_path)}")
                    except Exception as e:
                        QMessageBox.warning(self, "导入失败", str(e))
                    break
            event.acceptProposedAction()
        else:
            super().dropEvent(event)


class PPTGeneratorTool(QMainWindow):
    """主窗口"""

    def __init__(self):
        super().__init__()
        self.settings = QSettings("PPTGenerator", "OutlineToPPT")
        self.dark_mode = False
        self.template_path = None
        self._init_ui()
        self._init_menu()
        self._init_toolbar()
        self._init_statusbar()
        self._load_settings()
        self._apply_theme()

    def _init_ui(self):
        self.setWindowTitle("大纲转 PPT 工具 v1.1")
        self.resize(1050, 750)
        self.setMinimumSize(850, 600)

        central = QWidget()
        self.setCentralWidget(central)
        main_layout = QHBoxLayout(central)
        main_layout.setContentsMargins(10, 10, 10, 10)
        main_layout.setSpacing(10)

        # ===== 左侧：输入区 =====
        left = QWidget()
        left_layout = QVBoxLayout(left)
        left_layout.setContentsMargins(0, 0, 0, 0)

        input_group = QGroupBox("📝 大纲内容（支持拖拽 .md/.txt）")
        input_layout = QVBoxLayout(input_group)

        self.text_edit = DragDropTextEdit()
        self.text_edit.setPlaceholderText(
            "【示例】\n\n"
            "# 演示文稿标题\n"
            "副标题内容\n"
            "---\n"
            "## 第一章\n"
            "* 要点一\n"
            "* 要点二\n"
            "---\n"
            "## 第二章\n"
            "正文内容...\n"
        )
        self.text_edit.setFont(QFont("Consolas", 11))
        input_layout.addWidget(self.text_edit)

        self.char_label = QLabel("字符: 0 | 行: 0")
        self.char_label.setStyleSheet("color: #666;")
        self.text_edit.textChanged.connect(self._update_stats)
        input_layout.addWidget(self.char_label)

        left_layout.addWidget(input_group)

        # ===== 右侧：设置区 =====
        right = QWidget()
        right.setFixedWidth(320)
        right_layout = QVBoxLayout(right)
        right_layout.setContentsMargins(0, 0, 0, 0)
        right_layout.setSpacing(8)

        # 模板设置
        tpl_group = QGroupBox("📁 模板")
        tpl_layout = QHBoxLayout()
        self.template_label = QLabel("默认模板")
        self.template_label.setStyleSheet("color: #666;")
        tpl_layout.addWidget(self.template_label, 1)
        self.select_tpl_btn = QPushButton("选择")
        self.select_tpl_btn.setFixedWidth(60)
        self.select_tpl_btn.clicked.connect(self._select_template)
        tpl_layout.addWidget(self.select_tpl_btn)
        self.clear_tpl_btn = QPushButton("清除")
        self.clear_tpl_btn.setFixedWidth(50)
        self.clear_tpl_btn.clicked.connect(self._clear_template)
        tpl_layout.addWidget(self.clear_tpl_btn)
        tpl_group.setLayout(tpl_layout)
        right_layout.addWidget(tpl_group)

        # 基础设置
        basic_group = QGroupBox("⚙️ 基础设置")
        form1 = QFormLayout()
        form1.setSpacing(8)

        self.separator_input = QLineEdit("---")
        form1.addRow("分页符:", self.separator_input)

        self.font_combo = QComboBox()
        self.font_map = {
            "微软雅黑": "Microsoft YaHei",
            "黑体": "SimHei",
            "宋体": "SimSun",
            "仿宋": "FangSong",
            "楷体": "KaiTi",
        }
        self.font_combo.addItems(self.font_map.keys())
        form1.addRow("中文字体:", self.font_combo)

        self.latin_font_combo = QComboBox()
        self.latin_font_map = {
            "Times New Roman": "Times New Roman",
            "Arial": "Arial",
            "Calibri": "Calibri",
            "Consolas": "Consolas",
        }
        self.latin_font_combo.addItems(self.latin_font_map.keys())
        form1.addRow("英文/数字:", self.latin_font_combo)

        self.title_size_spin = QSpinBox()
        self.title_size_spin.setRange(16, 72)
        self.title_size_spin.setValue(32)
        self.title_size_spin.setSuffix(" pt")
        form1.addRow("标题字号:", self.title_size_spin)

        self.body_size_spin = QSpinBox()
        self.body_size_spin.setRange(10, 48)
        self.body_size_spin.setValue(20)
        self.body_size_spin.setSuffix(" pt")
        form1.addRow("正文字号:", self.body_size_spin)

        basic_group.setLayout(form1)
        right_layout.addWidget(basic_group)

        # 段落格式
        para_group = QGroupBox("📐 段落格式")
        form2 = QFormLayout()
        form2.setSpacing(8)

        self.indent_spin = QSpinBox()
        self.indent_spin.setRange(0, 8)
        self.indent_spin.setValue(2)
        self.indent_spin.setSuffix(" 字符")
        form2.addRow("首行缩进:", self.indent_spin)

        self.line_spacing_spin = QDoubleSpinBox()
        self.line_spacing_spin.setRange(1.0, 3.0)
        self.line_spacing_spin.setValue(1.5)
        self.line_spacing_spin.setSingleStep(0.1)
        self.line_spacing_spin.setSuffix(" 倍")
        form2.addRow("行距:", self.line_spacing_spin)

        self.para_spacing_spin = QSpinBox()
        self.para_spacing_spin.setRange(0, 30)
        self.para_spacing_spin.setValue(0)
        self.para_spacing_spin.setSuffix(" pt")
        form2.addRow("段前段后:", self.para_spacing_spin)

        para_group.setLayout(form2)
        right_layout.addWidget(para_group)

        # 配色
        style_group = QGroupBox("🎨 配色")
        form3 = QFormLayout()
        self.theme_combo = QComboBox()
        self.theme_combo.addItems(THEMES.keys())
        form3.addRow("方案:", self.theme_combo)
        style_group.setLayout(form3)
        right_layout.addWidget(style_group)

        # 选项
        opt_group = QGroupBox("🔧 选项")
        opt_layout = QVBoxLayout()
        opt_layout.setSpacing(5)

        self.clean_md_checkbox = QCheckBox("清理 Markdown 符号")
        self.clean_md_checkbox.setChecked(True)
        opt_layout.addWidget(self.clean_md_checkbox)

        self.cover_checkbox = QCheckBox("生成封面页")
        self.cover_checkbox.setChecked(True)
        opt_layout.addWidget(self.cover_checkbox)

        self.toc_checkbox = QCheckBox("生成目录页")
        self.toc_checkbox.setChecked(False)
        opt_layout.addWidget(self.toc_checkbox)

        self.open_after_checkbox = QCheckBox("导出后打开")
        self.open_after_checkbox.setChecked(True)
        opt_layout.addWidget(self.open_after_checkbox)

        opt_group.setLayout(opt_layout)
        right_layout.addWidget(opt_group)

        right_layout.addStretch()

        # 预览
        info_frame = QFrame()
        info_frame.setStyleSheet("background:#f0f0f0;border-radius:6px;")
        info_layout = QVBoxLayout(info_frame)
        info_layout.setContentsMargins(10, 8, 10, 8)
        self.preview_label = QLabel("📊 预计: 0 页")
        self.preview_label.setStyleSheet("font-weight:bold;color:#0078d4;")
        self.text_edit.textChanged.connect(self._update_preview)
        info_layout.addWidget(self.preview_label)
        right_layout.addWidget(info_frame)

        # 导出按钮
        self.export_btn = QPushButton("📤 生成 PPT")
        self.export_btn.setMinimumHeight(50)
        self.export_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self.export_btn.setStyleSheet("""
            QPushButton {
                background-color: #0078D4; color: white;
                font-size: 14px; font-weight: bold;
                border: none; border-radius: 8px;
            }
            QPushButton:hover { background-color: #1a86d9; }
            QPushButton:pressed { background-color: #005a9e; }
        """)
        self.export_btn.clicked.connect(self._on_export)
        right_layout.addWidget(self.export_btn)

        main_layout.addWidget(left, 7)
        main_layout.addWidget(right, 3)

    def _init_menu(self):
        menubar = self.menuBar()

        file_menu = menubar.addMenu("文件(&F)")
        open_act = QAction("打开(&O)", self)
        open_act.setShortcut(QKeySequence.StandardKey.Open)
        open_act.triggered.connect(self._open_file)
        file_menu.addAction(open_act)

        tpl_act = QAction("选择模板(&T)", self)
        tpl_act.setShortcut("Ctrl+T")
        tpl_act.triggered.connect(self._select_template)
        file_menu.addAction(tpl_act)

        save_act = QAction("导出(&S)", self)
        save_act.setShortcut(QKeySequence.StandardKey.Save)
        save_act.triggered.connect(self._on_export)
        file_menu.addAction(save_act)

        file_menu.addSeparator()
        exit_act = QAction("退出(&Q)", self)
        exit_act.setShortcut("Ctrl+Q")
        exit_act.triggered.connect(self.close)
        file_menu.addAction(exit_act)

        edit_menu = menubar.addMenu("编辑(&E)")
        clear_act = QAction("清空", self)
        clear_act.triggered.connect(lambda: self.text_edit.clear())
        edit_menu.addAction(clear_act)

        view_menu = menubar.addMenu("视图(&V)")
        self.dark_act = QAction("深色模式", self)
        self.dark_act.setCheckable(True)
        self.dark_act.triggered.connect(self._toggle_dark)
        view_menu.addAction(self.dark_act)

        help_menu = menubar.addMenu("帮助(&H)")
        about_act = QAction("关于(&A)", self)
        about_act.triggered.connect(self._show_about)
        help_menu.addAction(about_act)

    def _init_toolbar(self):
        tb = QToolBar()
        tb.setMovable(False)
        self.addToolBar(tb)

        tb.addAction("📂 打开", self._open_file)
        tb.addAction("📋 模板", self._select_template)
        tb.addAction("💾 导出", self._on_export)
        tb.addSeparator()
        tb.addAction("🗑️ 清空", lambda: self.text_edit.clear())

    def _init_statusbar(self):
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)
        self.status_bar.showMessage("就绪")

    def _select_template(self):
        path, _ = QFileDialog.getOpenFileName(self, "选择模板", "", "PowerPoint (*.pptx)")
        if path:
            self.template_path = path
            self.template_label.setText(os.path.basename(path))
            self.template_label.setStyleSheet("color:#0078d4;font-weight:bold;")
            self.status_bar.showMessage(f"模板: {path}")

    def _clear_template(self):
        self.template_path = None
        self.template_label.setText("默认模板")
        self.template_label.setStyleSheet("color:#666;")

    def _load_settings(self):
        try:
            self.font_combo.setCurrentText(self.settings.value("font", "微软雅黑"))
            self.latin_font_combo.setCurrentText(self.settings.value("latin_font", "Times New Roman"))
            self.title_size_spin.setValue(int(self.settings.value("title_size", 32)))
            self.body_size_spin.setValue(int(self.settings.value("body_size", 20)))
            self.indent_spin.setValue(int(self.settings.value("indent", 2)))
            self.line_spacing_spin.setValue(float(self.settings.value("line_spacing", 1.5)))
            self.para_spacing_spin.setValue(int(self.settings.value("para_spacing", 0)))
            self.theme_combo.setCurrentText(self.settings.value("theme", "经典蓝"))
            self.cover_checkbox.setChecked(self.settings.value("cover", True, type=bool))
            self.toc_checkbox.setChecked(self.settings.value("toc", False, type=bool))
            self.dark_mode = self.settings.value("dark_mode", False, type=bool)
            self.dark_act.setChecked(self.dark_mode)
            tpl = self.settings.value("template_path", "")
            if tpl and os.path.exists(tpl):
                self.template_path = tpl
                self.template_label.setText(os.path.basename(tpl))
                self.template_label.setStyleSheet("color:#0078d4;font-weight:bold;")
        except:
            pass

    def _save_settings(self):
        try:
            self.settings.setValue("font", self.font_combo.currentText())
            self.settings.setValue("latin_font", self.latin_font_combo.currentText())
            self.settings.setValue("title_size", self.title_size_spin.value())
            self.settings.setValue("body_size", self.body_size_spin.value())
            self.settings.setValue("indent", self.indent_spin.value())
            self.settings.setValue("line_spacing", self.line_spacing_spin.value())
            self.settings.setValue("para_spacing", self.para_spacing_spin.value())
            self.settings.setValue("theme", self.theme_combo.currentText())
            self.settings.setValue("cover", self.cover_checkbox.isChecked())
            self.settings.setValue("toc", self.toc_checkbox.isChecked())
            self.settings.setValue("dark_mode", self.dark_mode)
            self.settings.setValue("template_path", self.template_path or "")
        except:
            pass

    def _toggle_dark(self):
        self.dark_mode = not self.dark_mode
        self._apply_theme()

    def _apply_theme(self):
        if self.dark_mode:
            self.setStyleSheet("""
                QMainWindow, QWidget { background-color: #2b2b2b; color: #e0e0e0; }
                QGroupBox { border: 1px solid #555; border-radius: 6px; margin-top: 10px; padding-top: 10px; }
                QTextEdit, QLineEdit, QSpinBox, QDoubleSpinBox, QComboBox {
                    background-color: #3c3c3c; color: #e0e0e0;
                    border: 1px solid #555; border-radius: 4px; padding: 5px;
                }
                QMenuBar { background-color: #2b2b2b; color: #e0e0e0; }
                QMenu { background-color: #2b2b2b; color: #e0e0e0; border: 1px solid #555; }
                QMenu::item:selected { background-color: #0078d4; }
                QToolBar, QStatusBar { background-color: #2b2b2b; border: none; }
                QFrame { background-color: #3c3c3c; }
                QPushButton { background-color: #3c3c3c; color: #e0e0e0; border: 1px solid #555; border-radius: 4px; padding: 5px; }
            """)
            self.preview_label.setStyleSheet("font-weight:bold;color:#4da6ff;")
        else:
            self.setStyleSheet("")
            self.preview_label.setStyleSheet("font-weight:bold;color:#0078d4;")

    def _update_stats(self):
        text = self.text_edit.toPlainText()
        self.char_label.setText(f"字符: {len(text)} | 行: {text.count(chr(10)) + 1 if text else 0}")

    def _update_preview(self):
        text = self.text_edit.toPlainText().strip()
        sep = self.separator_input.text() or "---"
        if not text:
            self.preview_label.setText("📊 预计: 0 页")
            return
        blocks = [b.strip() for b in text.split(sep) if b.strip()]
        n = len(blocks)
        extra = ""
        if self.cover_checkbox.isChecked() and n > 0:
            extra = "(含封面)"
        if self.toc_checkbox.isChecked() and n > 1:
            n += 1
            extra += "+目录"
        self.preview_label.setText(f"📊 预计: {n} 页 {extra}")

    def _open_file(self):
        path, _ = QFileDialog.getOpenFileName(self, "打开", "", "文本 (*.txt *.md);;所有 (*.*)")
        if path:
            try:
                content = None
                for enc in ['utf-8', 'gbk', 'gb2312']:
                    try:
                        with open(path, 'r', encoding=enc) as f:
                            content = f.read()
                        break
                    except:
                        continue
                if content:
                    self.text_edit.setPlainText(content)
                    self.status_bar.showMessage(f"已打开: {path}")
            except Exception as e:
                QMessageBox.warning(self, "失败", str(e))

    def _clean_markdown(self, text: str) -> str:
        """彻底清理 Markdown"""
        # 标题 # 符号
        text = re.sub(r'^[ \t]*#{1,6}[ \t]+', '', text, flags=re.MULTILINE)
        text = re.sub(r'^[ \t]*#{1,6}[ \t]*$', '', text, flags=re.MULTILINE)
        
        # 加粗斜体
        text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'___(.+?)___', r'\1', text)
        text = re.sub(r'__(.+?)__', r'\1', text)
        text = re.sub(r'(?<![*])\*([^*\n]+?)\*(?![*])', r'\1', text)
        text = re.sub(r'(?<![_])_([^_\n]+?)_(?![_])', r'\1', text)
        
        # 删除线、代码
        text = re.sub(r'~~(.+?)~~', r'\1', text)
        text = re.sub(r'`([^`\n]+?)`', r'\1', text)
        
        # 链接、图片
        text = re.sub(r'\[([^\]]+?)\]\([^)]+?\)', r'\1', text)
        text = re.sub(r'!\[([^\]]*?)\]\([^)]+?\)', r'\1', text)
        
        # 列表符号
        text = re.sub(r'^[ \t]*[\*\-\+][ \t]+', '• ', text, flags=re.MULTILINE)
        text = re.sub(r'^[ \t]*\d+\.[ \t]+', '', text, flags=re.MULTILINE)
        
        # 引用、代码块
        text = re.sub(r'^[ \t]*>+[ \t]*', '', text, flags=re.MULTILINE)
        text = re.sub(r'^```.*$', '', text, flags=re.MULTILINE)
        
        # 多余空行
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text

    def _set_run_font(self, run, cn_font, latin_font, size, color=None, bold=False):
        """设置文字样式"""
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = get_rgb_color(color)

        try:
            rPr = run._r.get_or_add_rPr()
            
            # 拉丁字体
            latin = rPr.find(qn('a:latin'))
            if latin is None:
                latin = etree.SubElement(rPr, qn('a:latin'))
            latin.set('typeface', latin_font)
            
            # 东亚字体
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', cn_font)
            
            # 复杂脚本
            cs = rPr.find(qn('a:cs'))
            if cs is None:
                cs = etree.SubElement(rPr, qn('a:cs'))
            cs.set('typeface', latin_font)
        except:
            run.font.name = cn_font

    def _set_paragraph_format(self, para, font_size, indent_chars=0, line_spacing=1.5, 
                               space_before=0, space_after=0, is_title=False):
        """
        设置段落格式 (修复版)
        使用 python-pptx 正确的属性和 XML 操作
        """
        # 段前段后 (直接设置)
        para.space_before = Pt(space_before)
        para.space_after = Pt(space_after)
        
        # 行距 (直接设置倍数)
        para.line_spacing = line_spacing

        # 首行缩进 (通过 XML 设置)
        if not is_title and indent_chars > 0:
            try:
                # 获取段落的 XML 元素
                pPr = para._p.get_or_add_pPr()
                # 计算缩进值 (EMU)
                indent_emu = int(Pt(indent_chars * font_size))
                # 设置 indent 属性
                pPr.set('indent', str(indent_emu))
            except Exception as e:
                print(f"缩进设置警告: {e}")

    def _on_export(self):
        content = self.text_edit.toPlainText().strip()
        if not content:
            QMessageBox.warning(self, "提示", "请先输入内容！")
            return

        path, _ = QFileDialog.getSaveFileName(self, "保存", "演示文稿.pptx", "PowerPoint (*.pptx)")
        if not path:
            return
        if not path.lower().endswith('.pptx'):
            path += '.pptx'

        self._save_settings()

        try:
            count = self._generate_ppt(content, path)
            msg = f"成功生成 {count} 页！\n\n{path}"
            
            if self.open_after_checkbox.isChecked():
                reply = QMessageBox.information(
                    self, "成功 ✓", msg,
                    QMessageBox.StandardButton.Open | QMessageBox.StandardButton.Ok,
                    QMessageBox.StandardButton.Open
                )
                if reply == QMessageBox.StandardButton.Open:
                    self._open_external(path)
            else:
                QMessageBox.information(self, "成功 ✓", msg)
            
            self.status_bar.showMessage(f"已导出: {path}")
        except PermissionError:
            QMessageBox.critical(self, "失败", "文件被占用，请关闭后重试！")
        except Exception as e:
            QMessageBox.critical(self, "失败", f"错误: {e}")
            import traceback
            traceback.print_exc()

    def _open_external(self, path):
        try:
            if sys.platform == 'win32':
                os.startfile(path)
            elif sys.platform == 'darwin':
                subprocess.call(['open', path])
            else:
                subprocess.call(['xdg-open', path])
        except:
            pass

    def _generate_ppt(self, text: str, output_path: str) -> int:
        """生成 PPT"""
        sep = self.separator_input.text() or "---"
        cn_font = self.font_map.get(self.font_combo.currentText(), "Microsoft YaHei")
        latin_font = self.latin_font_map.get(self.latin_font_combo.currentText(), "Times New Roman")
        title_size = self.title_size_spin.value()
        body_size = self.body_size_spin.value()
        indent = self.indent_spin.value()
        line_sp = self.line_spacing_spin.value()
        para_sp = self.para_spacing_spin.value()
        theme = THEMES.get(self.theme_combo.currentText(), THEMES["经典蓝"])
        clean_md = self.clean_md_checkbox.isChecked()
        make_cover = self.cover_checkbox.isChecked()
        make_toc = self.toc_checkbox.isChecked()

        # 创建 PPT
        if self.template_path and os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        blocks = [b.strip() for b in text.split(sep) if b.strip()]
        if not blocks:
            raise ValueError("无有效内容")

        slide_count = 0
        toc_titles = []

        # ===== 封面页 =====
        if make_cover and blocks:
            block = blocks[0]
            if clean_md:
                block = self._clean_markdown(block)
            lines = [l.strip() for l in block.splitlines() if l.strip()]

            slide = prs.slides.add_slide(prs.slide_layouts[0])

            if lines and slide.shapes.title:
                slide.shapes.title.text = lines[0]
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    self._set_paragraph_format(p, title_size + 8, 0, 1.2, 0, 0, True)
                    for r in p.runs:
                        self._set_run_font(r, cn_font, latin_font, title_size + 8, theme["title_color"], True)

            if len(lines) > 1 and len(slide.placeholders) > 1:
                sub = slide.placeholders[1]
                sub.text = "\n".join(lines[1:])
                for p in sub.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    self._set_paragraph_format(p, body_size, 0, 1.5, 0, 0, True)
                    for r in p.runs:
                        self._set_run_font(r, cn_font, latin_font, body_size, theme["body_color"])

            blocks = blocks[1:]
            slide_count += 1

        # 收集目录
        for block in blocks:
            tmp = self._clean_markdown(block) if clean_md else block
            lines = [l.strip() for l in tmp.splitlines() if l.strip()]
            if lines:
                toc_titles.append(lines[0])

        # ===== 目录页 =====
        if make_toc and toc_titles:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            if slide.shapes.title:
                slide.shapes.title.text = "目录"
                for p in slide.shapes.title.text_frame.paragraphs:
                    self._set_paragraph_format(p, title_size, 0, 1.2, 0, 0, True)
                    for r in p.runs:
                        self._set_run_font(r, cn_font, latin_font, title_size, theme["title_color"], True)

            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, title in enumerate(toc_titles):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"{i + 1}. {title}"
                    p.level = 0
                    self._set_paragraph_format(p, body_size, 0, line_sp, para_sp, para_sp)
                    for r in p.runs:
                        self._set_run_font(r, cn_font, latin_font, body_size, theme["body_color"], True)

            slide_count += 1

        # ===== 内容页 =====
        for block in blocks:
            if clean_md:
                block = self._clean_markdown(block)

            lines = [l for l in block.splitlines() if l.strip()]
            if not lines:
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[1])

            # 标题
            title_text = lines[0].strip()
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    self._set_paragraph_format(p, title_size, 0, 1.2, 0, 0, True)
                    for r in p.runs:
                        self._set_run_font(r, cn_font, latin_font, title_size, theme["title_color"], True)

            # 正文
            body_lines = lines[1:]
            if body_lines and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()

                first = True
                for line in body_lines:
                    orig = line
                    line_stripped = line.strip()
                    if not line_stripped:
                        continue

                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = line_stripped

                    # 缩进层级
                    level = 0
                    tmp = orig
                    while tmp.startswith('\t') or tmp.startswith('    '):
                        level += 1
                        tmp = tmp[1:] if tmp.startswith('\t') else tmp[4:]
                    p.level = min(level, 4)

                    # 段落格式
                    self._set_paragraph_format(p, body_size, indent, line_sp, para_sp, para_sp)

                    # 字体
                    for r in p.runs:
                        self._set_run_font(r, cn_font, latin_font, body_size, theme["body_color"])

            slide_count += 1

        prs.save(output_path)
        return slide_count

    def _show_about(self):
        QMessageBox.about(
            self, "关于",
            "<h3>大纲转PPT v1.1</h3>"
            "<p>Markdown/文本 → PowerPoint</p>"
            "<hr><b>功能:</b><ul>"
            "<li>自定义模板</li>"
            "<li>中英文字体分设</li>"
            "<li>首行缩进、行距、段距</li>"
            "<li>Markdown 清理</li>"
            "<li>封面页+目录页</li>"
            "</ul>"
        )

    def closeEvent(self, event):
        self._save_settings()
        super().closeEvent(event)


def main():
    app = QApplication(sys.argv)
    app.setStyle("Fusion")
    app.setFont(QFont("Microsoft YaHei", 9))
    win = PPTGeneratorTool()
    win.show()
    sys.exit(app.exec())


if __name__ == "__main__":

    main()
